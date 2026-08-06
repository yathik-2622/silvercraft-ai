"""
document_text_extract — turns an uploaded KB reference file into plain
text. Unlike app/tools/upload_ingestion.py (which deliberately extracts
ONLY structural/aggregate metadata and never the content, per TDS §3),
this one is for reference/knowledge documents, where the actual text IS
the thing we need — that's the point of a Modelling Reference KB. This
distinction matters: source data files (customer records etc.) are never
stored; KB reference documents (modeling methodology docs an admin
explicitly uploads to teach the system) are stored in full, deliberately,
because citations need the whole document to render a highlight against.

Supports .md/.markdown/.txt, .pdf, .docx, .pptx, .csv, .xlsx.
PDF/PPTX extraction joins pages/slides with a form-feed (\\f) so the
"page_aware" chunking strategy (app/core/chunking.py) can detect and
respect real page/slide boundaries.

Parser choice per format (mirrors hckb_core_parsers.py's per-extension
loader-dispatch shape, but not a blind 1:1 swap to LangChain everywhere —
see the docstring on each function below for why):
  .pdf  -> LangChain's PyPDFLoader (same underlying pypdf, but Document-
           per-page output is what a LangChain-based pipeline expects).
  .csv  -> LangChain's CSVLoader (new capability — CSV reference docs,
           e.g. a data dictionary, weren't ingestable as KB text before).
  .docx -> kept on native python-docx, NOT LangChain's Docx2txtLoader —
           verified empirically that Docx2txtLoader flattens table rows
           into one cell per line with no row grouping, silently losing
           the "which cell belongs to which row" structure a modeling
           rules table needs; python-docx's " | "-per-row join preserves
           it and was worth keeping.
  .pptx -> kept on native python-pptx — LangChain's only PPTX loader
           (UnstructuredPowerPointLoader) requires the `unstructured`
           package, which pulls in ~40 transitive dependencies including
           spaCy/numba/llvmlite (60MB+) for a format this app doesn't
           receive often; not a good size/complexity trade for equivalent
           output.
  .xlsx -> openpyxl directly (same library app/tools/upload_ingestion.py
           already uses) — LangChain's UnstructuredExcelLoader has the
           same `unstructured` dependency cost as the PPTX loader above.
"""
import logging
import os
import tempfile
from typing import BinaryIO

import docx
import openpyxl
import pptx
from langchain_community.document_loaders import CSVLoader, PyPDFLoader

logger = logging.getLogger(__name__)

ADM_SUPPORTED_KB_EXTENSIONS = {".md", ".markdown", ".txt", ".pdf", ".docx", ".pptx", ".csv", ".xlsx"}


def ADM_extract_document_text(file_obj: BinaryIO, filename: str) -> str:
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ADM_SUPPORTED_KB_EXTENSIONS:
        raise ValueError(f"Unsupported KB document type '{ext}'. Supported: {sorted(ADM_SUPPORTED_KB_EXTENSIONS)}")

    logger.info("Extracting KB document text: filename=%r ext=%s", filename, ext)
    if ext == ".pdf":
        return ADM_extract_pdf_text(file_obj)
    if ext == ".docx":
        return ADM_extract_docx_text(file_obj)
    if ext == ".pptx":
        return ADM_extract_pptx_text(file_obj)
    if ext == ".csv":
        return ADM_extract_csv_text(file_obj)
    if ext == ".xlsx":
        return ADM_extract_xlsx_text(file_obj)

    raw = file_obj.read()
    return raw.decode("utf-8", errors="replace")


def _load_via_temp_file(file_obj: BinaryIO, suffix: str, loader_factory):
    """LangChain loaders read from a file PATH, not a file-like object
    (hckb_core_parsers.py does the exact same temp-file staging, for the
    same reason) — write the bytes out, always clean up in `finally`
    even if the loader raises."""
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_obj.read())
        tmp_path = tmp.name
    try:
        return loader_factory(tmp_path).load()
    finally:
        os.unlink(tmp_path)


def ADM_extract_pdf_text(file_obj: BinaryIO) -> str:
    """PyPDFLoader's default mode="page" returns one Document per page —
    join with \\f ourselves, identical contract to the previous pypdf-based
    per-page join."""
    docs = _load_via_temp_file(file_obj, ".pdf", lambda path: PyPDFLoader(path))
    return "\f".join(d.page_content for d in docs)


def ADM_extract_docx_text(file_obj: BinaryIO) -> str:
    """
    Paragraphs AND table content — a Word doc's tables often carry the
    actual modeling rules (e.g. a normalization-rules table), so skipping
    them would silently drop real content. Kept on native python-docx —
    see module docstring for why this one isn't routed through LangChain.
    """
    document = docx.Document(file_obj)
    parts: list[str] = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                parts.append(row_text)
    return "\n\n".join(parts)


def ADM_extract_pptx_text(file_obj: BinaryIO) -> str:
    """One slide's shapes joined per slide, slides joined by form-feed so
    page_aware chunking can treat each slide as one page."""
    presentation = pptx.Presentation(file_obj)
    slides_text: list[str] = []
    for slide in presentation.slides:
        shape_texts = [shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
        slides_text.append("\n".join(shape_texts))
    return "\f".join(slides_text)


def ADM_extract_csv_text(file_obj: BinaryIO) -> str:
    """CSVLoader returns one Document per row (page_content formatted as
    "col: value" lines per row) — good, embeddable reference text for
    something like a data dictionary or a standards spreadsheet. Rows
    joined with a blank line each, same "one flat string, chunked later"
    contract every other extractor here returns."""
    docs = _load_via_temp_file(file_obj, ".csv", lambda path: CSVLoader(path))
    return "\n\n".join(d.page_content for d in docs)


def ADM_extract_xlsx_text(file_obj: BinaryIO) -> str:
    """openpyxl directly (same library app/tools/upload_ingestion.py's
    profiling path already uses) — see module docstring for why this one
    isn't routed through LangChain's UnstructuredExcelLoader. Sheets
    joined by form-feed (page_aware chunking treats each sheet as a
    "page"); within a sheet, one line per row, cells comma-joined."""
    workbook = openpyxl.load_workbook(file_obj, read_only=True, data_only=True)
    sheets_text: list[str] = []
    for sheet in workbook.worksheets:
        lines = [
            ", ".join("" if c is None else str(c) for c in row)
            for row in sheet.iter_rows(values_only=True)
            if any(c is not None for c in row)
        ]
        if lines:
            sheets_text.append(f"[Sheet: {sheet.title}]\n" + "\n".join(lines))
    return "\f".join(sheets_text)
