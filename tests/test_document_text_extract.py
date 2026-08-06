"""
Unit tests for app/tools/document_text_extract.py after the Phase 3
LangChain swap (PDF -> PyPDFLoader, new CSV -> CSVLoader, new XLSX ->
openpyxl, DOCX/PPTX kept native — see the module docstring for why).
Run with: pytest

All fixtures are built in-memory (no external files) — real small
documents constructed with the same libraries the app itself uses to
write them (pypdf's PdfWriter, python-docx, python-pptx, openpyxl), so
these are genuine round-trips, not mocks.
"""
import io

import docx
import openpyxl
import pptx
from pypdf import PdfWriter

from app.tools.document_text_extract import ADM_extract_document_text, ADM_SUPPORTED_KB_EXTENSIONS


def test_unsupported_extension_raises_with_clear_message():
    try:
        ADM_extract_document_text(io.BytesIO(b"whatever"), "file.json")
        assert False, "expected ValueError"
    except ValueError as e:
        assert "Unsupported KB document type '.json'" in str(e)
        assert "Supported:" in str(e)


def test_txt_and_md_still_plain_decode():
    assert ADM_extract_document_text(io.BytesIO(b"hello world"), "notes.txt") == "hello world"
    assert ADM_extract_document_text(io.BytesIO(b"# Heading"), "notes.md") == "# Heading"


def test_pdf_pages_join_with_form_feed():
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.add_blank_page(width=200, height=200)
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)

    text = ADM_extract_document_text(buf, "doc.pdf")
    # page_aware chunking (app/core/chunking.py) splits on \f — the exact
    # contract this test guards even though blank pages have no real text.
    assert len(text.split("\f")) == 3


def test_docx_preserves_paragraphs_and_table_rows():
    document = docx.Document()
    document.add_paragraph("A modeling rule paragraph.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Column"
    table.cell(0, 1).text = "Rule"
    table.cell(1, 0).text = "customer_id"
    table.cell(1, 1).text = "Must be unique"
    buf = io.BytesIO()
    document.save(buf)
    buf.seek(0)

    text = ADM_extract_document_text(buf, "rules.docx")
    assert "A modeling rule paragraph." in text
    # Row structure preserved (this is exactly what Docx2txtLoader was
    # found NOT to do — see module docstring) — a whole row survives as
    # one " | "-joined line, not scattered across separate lines.
    assert "Column | Rule" in text
    assert "customer_id | Must be unique" in text


def test_pptx_slides_join_with_form_feed():
    presentation = pptx.Presentation()
    for _ in range(2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(0, 0, 100, 100)
        box.text_frame.text = "Slide text"
    buf = io.BytesIO()
    presentation.save(buf)
    buf.seek(0)

    text = ADM_extract_document_text(buf, "deck.pptx")
    slides = text.split("\f")
    assert len(slides) == 2
    assert all("Slide text" in s for s in slides)


def test_csv_extracts_row_content():
    csv_bytes = b"name,role\nAlice,Engineer\nBob,Analyst\n"
    text = ADM_extract_document_text(io.BytesIO(csv_bytes), "people.csv")
    assert "Alice" in text
    assert "Engineer" in text
    assert "Bob" in text


def test_xlsx_extracts_sheet_content_with_form_feed_per_sheet():
    workbook = openpyxl.Workbook()
    ws1 = workbook.active
    ws1.title = "Sheet1"
    ws1.append(["col1", "col2"])
    ws1.append(["a", 1])
    ws2 = workbook.create_sheet("Sheet2")
    ws2.append(["x", "y"])
    buf = io.BytesIO()
    workbook.save(buf)
    buf.seek(0)

    text = ADM_extract_document_text(buf, "data.xlsx")
    sheets = text.split("\f")
    assert len(sheets) == 2
    assert "[Sheet: Sheet1]" in sheets[0]
    assert "col1, col2" in sheets[0]
    assert "[Sheet: Sheet2]" in sheets[1]


def test_supported_extensions_set_is_a_superset_of_pre_phase3():
    pre_phase3 = {".md", ".markdown", ".txt", ".pdf", ".docx", ".pptx"}
    assert pre_phase3.issubset(ADM_SUPPORTED_KB_EXTENSIONS)
    assert {".csv", ".xlsx"}.issubset(ADM_SUPPORTED_KB_EXTENSIONS)
